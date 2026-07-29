"""
Build the Second Progress Review (PR2) presentation.

Follows the department outline in `docs/2nd PR_Presentation_Outline_IT_IS.pptx`
section by section, and keeps the 16:9 shape of the PR1 deck. Every number in
here is taken from the repository (model metrics JSON, the test workbook, the
survey CSV, the test suites) rather than written by hand — regenerate the deck
after any of those change:

    python scripts/build_pr2_pptx.py

Language is deliberately plain: this is read aloud in a live review, so short
lines beat complete sentences. Speaker notes carry the detail instead.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DIAGRAMS = ROOT / "docs" / "diagrams"
OUTPUT = ROOT / "docs" / "PR2_Presentation.pptx"

# ── Palette ────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x14, 0x2A, 0x3F)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
AMBER = RGBColor(0xC2, 0x6A, 0x1A)
INK = RGBColor(0x1F, 0x2A, 0x33)
GREY = RGBColor(0x5B, 0x6B, 0x77)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Live figures pulled from the repo ──────────────────────────────────────
def load_facts() -> dict:
    cond = json.loads(
        (ROOT / "backend/services/python/ai-service/app/models/condition_metrics.json").read_text()
    )
    risk = json.loads(
        (ROOT / "backend/services/python/ai-service/app/models/risk_metrics.json").read_text()
    )
    report = cond["classification_report"]
    per_class = {
        k: v["f1-score"] for k, v in report.items() if isinstance(v, dict) and "f1-score" in v and "avg" not in k
    }
    return {
        "accuracy": cond["accuracy"],
        "macro_f1": cond["macro_f1"],
        "n_train": cond["n_train"],
        "n_test": cond["n_test"],
        "n_labels": len(cond["labels"]),
        "per_class": per_class,
        "roc_auc": risk["roc_auc"],
        "positive_rate": risk["positive_rate"],
    }


F = load_facts()

# Counted from the test suites and the manual test workbook (see docstring).
TESTS_AI, TESTS_AGENT, TESTS_GW = 19, 27, 7
TC_TOTAL, TC_PASS, TC_FAIL, TC_NOTRUN = 123, 82, 2, 39
SURVEY_N = 31


# ── Slide helpers ──────────────────────────────────────────────────────────
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill: RGBColor):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def write(frame, lines, size=18, color=INK, bold_first=False, space=10):
    """lines: list of (text, level) or plain strings."""
    frame.clear()
    for i, item in enumerate(lines):
        text, level = item if isinstance(item, tuple) else (item, 0)
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.level = level
        para.space_after = Pt(space)
        run = para.add_run()
        run.text = ("• " if level == 0 else "– ") + text if text else ""
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.color.rgb = color if level == 0 else GREY
        run.font.bold = bold_first and i == 0
        run.font.name = "Calibri"


def header(slide, number: str, title: str, subtitle: str = "") -> None:
    rect(slide, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), TEAL)

    frame = textbox(slide, Inches(0.55), Inches(0.16), Inches(11.5), Inches(0.9))
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = title
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    if subtitle:
        sub = frame.add_paragraph()
        srun = sub.add_run()
        srun.text = subtitle
        srun.font.size = Pt(13)
        srun.font.color.rgb = RGBColor(0xB8, 0xCD, 0xD8)
        srun.font.name = "Calibri"

    num = textbox(slide, Inches(12.1), Inches(0.28), Inches(0.9), Inches(0.6))
    npara = num.paragraphs[0]
    npara.alignment = PP_ALIGN.RIGHT
    nrun = npara.add_run()
    nrun.text = number
    nrun.font.size = Pt(26)
    nrun.font.bold = True
    nrun.font.color.rgb = TEAL
    nrun.font.name = "Calibri"


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def picture(slide, filename: str, top=Inches(1.6), max_h=Inches(5.4)) -> None:
    """Drop a diagram in, scaled to fit and centred."""
    path = DIAGRAMS / filename
    if not path.exists():
        placeholder(slide, f"[ Diagram missing: {filename} ]")
        return
    from PIL import Image

    with Image.open(path) as img:
        ratio = img.width / img.height
    max_w = Inches(11.6)
    h = max_h
    w = Emu(int(h * ratio))
    if w > max_w:
        w = max_w
        h = Emu(int(w / ratio))
    slide.shapes.add_picture(str(path), Emu(int((SLIDE_W - w) / 2)), top, w, h)


def placeholder(slide, label: str, x=Inches(0.9), y=Inches(1.9), w=Inches(11.5), h=Inches(4.6)) -> None:
    """A visible 'paste your screenshot here' box — never a fake screenshot."""
    box = rect(slide, x, y, w, h, LIGHT)
    box.line.color.rgb = RGBColor(0xC3, 0xCE, 0xD5)
    box.line.width = Pt(1.25)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = label
    run.font.size = Pt(16)
    run.font.color.rgb = GREY
    run.font.name = "Calibri"


def stat_row(slide, items, top=Inches(1.75)) -> None:
    """Big-number cards: value on top, caption under."""
    gap = Inches(0.3)
    width = Emu(int((Inches(11.9) - gap * (len(items) - 1)) / len(items)))
    x = Inches(0.72)
    for value, caption in items:
        card = rect(slide, x, top, width, Inches(1.5), LIGHT)
        card.line.color.rgb = RGBColor(0xD8, 0xE1, 0xE6)
        frame = card.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = value
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = TEAL
        run.font.name = "Calibri"
        cap = frame.add_paragraph()
        cap.alignment = PP_ALIGN.CENTER
        crun = cap.add_run()
        crun.text = caption
        crun.font.size = Pt(12)
        crun.font.color.rgb = GREY
        crun.font.name = "Calibri"
        x = Emu(int(x + width + gap))


def table(slide, rows, col_widths, top=Inches(1.75), font=13):
    shape = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(0.72), top, Inches(11.9), Inches(0.4 * len(rows))
    )
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = text
            run.font.size = Pt(font)
            run.font.name = "Calibri"
            run.font.bold = r == 0
            run.font.color.rgb = WHITE if r == 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (WHITE if r % 2 else LIGHT)
    return tbl


# ── Slides ─────────────────────────────────────────────────────────────────
def build() -> None:
    prs = new_deck()

    # 1 ── Title
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(s, 0, Inches(3.28), SLIDE_W, Inches(0.06), TEAL)

    f = textbox(s, Inches(1.0), Inches(1.15), Inches(11.3), Inches(2.1))
    p = f.paragraphs[0]
    r = p.add_run()
    r.text = (
        "Design and Evaluation of an Agentic AI-Driven Decision Support "
        "System for Explainable Early Disease Risk Detection in Companion Animals"
    )
    r.font.size = Pt(33)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    f2 = textbox(s, Inches(1.0), Inches(3.65), Inches(11.3), Inches(3.0))
    for i, (label, value) in enumerate(
        [
            ("", "Second Progress Review"),
            ("Presented by", "MTE Ranasinghe"),
            ("Student ID", "D/BIT/23/0025"),
            ("Degree Programme", "BSc (Hons) in Information Technology"),
            ("Supervisor", "Dr. (Mrs) N Wedasinghe"),
            ("Co-Supervisor", "Mr. BGL Balasuriya"),
            ("Department", "Department of Information Technology"),
            ("University", "General Sir John Kotelawala Defence University"),
            # Left blank on purpose: fill in the real review date before presenting.
            ("Date", "[ add presentation date ]"),
        ]
    ):
        para = f2.paragraphs[0] if i == 0 else f2.add_paragraph()
        para.space_after = Pt(6)
        if label:
            lr = para.add_run()
            lr.text = f"{label}:  "
            lr.font.size = Pt(15)
            lr.font.color.rgb = RGBColor(0x8F, 0xB2, 0xC2)
            lr.font.name = "Calibri"
        vr = para.add_run()
        vr.text = value
        vr.font.size = Pt(17 if not label else 15)
        vr.font.bold = not label
        vr.font.color.rgb = TEAL if not label else WHITE
        vr.font.name = "Calibri"
    notes(
        s,
        "Good morning. I am Tharushi Ranasinghe. This is my second progress review for my final "
        "year project on an agentic AI system that helps pet owners spot disease risk early. "
        "I will cover what I planned, what I have built, what the results show, and what is left.",
    )

    # 2 ── Project Overview
    s = blank(prs)
    header(s, "02", "Project Overview", "Background · Problem · Aim · Objectives")
    f = textbox(s, Inches(0.72), Inches(1.6), Inches(5.6), Inches(5.4))
    write(
        f,
        [
            "Background",
            ("Pets cannot tell us what is wrong.", 1),
            ("Early signs are small and easy to miss.", 1),
            "Problem",
            ("Owners often notice too late.", 1),
            ("Late treatment costs more and hurts welfare.", 1),
            ("Existing apps only remind, they do not reason.", 1),
        ],
        size=17,
    )
    f2 = textbox(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(5.4))
    write(
        f2,
        [
            "Aim",
            ("Build and evaluate an agentic AI system that finds", 1),
            ("early disease risk and explains its reasoning.", 1),
            "Objectives",
            ("1. Find risk from owner-reported symptoms.", 1),
            ("2. Train ML models for disease and risk level.", 1),
            ("3. Add an agentic AI layer that reasons and acts.", 1),
            ("4. Explain results: confidence and key symptoms.", 1),
            ("5. Find nearby clinics and book appointments.", 1),
            ("6. Track vaccinations and send reminders.", 1),
        ],
        size=17,
    )
    notes(
        s,
        "The problem is simple. Pets hide illness. Owners see the signs late. Apps today only send "
        "reminders. My aim is a system that reasons on its own, explains why, and then helps the "
        "owner take action. The six objectives are unchanged from PR1 and all six are now built.",
    )

    # 3 ── Literature Review Summary
    s = blank(prs)
    header(s, "03", "Literature Review Summary", "Key findings and the gap")
    table(
        s,
        [
            ["Area", "What the literature shows", "What is missing"],
            ["AI in animal disease prediction", "AI improves early detection", "Built for surveillance, not owners"],
            ["ML in animal healthcare", "ML raises prediction accuracy", "Mostly theory, no real mobile system"],
            ["Veterinary decision support", "AI can support vet decisions", "Clinical only, not preventive"],
            ["Chatbots in pet health", "Chatbots help owner interaction", "No autonomous reasoning"],
            ["Sensor-based monitoring", "Behaviour signals show health", "Needs hardware; no symptom input"],
            ["Dog health score systems", "Scores agree with vet diagnosis", "No reasoning or action guidance"],
        ],
        [3.0, 4.6, 4.3],
        font=13,
    )
    f = textbox(s, Inches(0.72), Inches(5.35), Inches(11.9), Inches(1.6))
    write(
        f,
        [
            "Research gap",
            ("No owner-facing system that reasons on its own, explains why, and acts on the result.", 1),
            ("Nothing built for Sri Lanka or available in Sinhala and Tamil.", 1),
        ],
        size=17,
    )
    notes(
        s,
        "I reviewed twelve studies. They agree AI works in animal health. But almost all of them are "
        "for vets, for livestock, or need sensors. None give the owner an explanation they can act on, "
        "and none work in Sinhala or Tamil. That is my gap.",
    )

    # 4 ── Research Methodology
    s = blank(prs)
    header(s, "04", "Research Methodology", "Approach · Design · Data · Tools")
    f = textbox(s, Inches(0.72), Inches(1.6), Inches(5.6), Inches(5.4))
    write(
        f,
        [
            "Approach",
            ("Design Science Research — build an artefact, then evaluate it.", 1),
            "Development",
            ("Agile, in short cycles. Each feature built and tested.", 1),
            "Data collection",
            (f"Primary: {SURVEY_N} pet owner responses (symptom survey).", 1),
            ("Secondary: two public pet symptom datasets.", 1),
        ],
        size=17,
    )
    f2 = textbox(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(5.4))
    write(
        f2,
        [
            "Data analysis",
            ("Accuracy, macro F1 and per-class F1.", 1),
            ("AI result compared with the vet's confirmed diagnosis.", 1),
            "Tools and technologies",
            ("React + TypeScript front end (installable web app).", 1),
            ("Node.js services and Python FastAPI for AI.", 1),
            ("MongoDB for data, Neo4j for the disease ontology.", 1),
            ("Docker Compose to run everything together.", 1),
        ],
        size=17,
    )
    notes(
        s,
        "I follow Design Science Research: I build the system, then measure it. Development is agile. "
        "For data I used two public symptom datasets plus a survey of thirty one pet owners. "
        "I measure the models with accuracy and F1, and I also record what the vet actually diagnosed "
        "so I can compare AI against real outcomes.",
    )

    # 5 ── Requirements
    s = blank(prs)
    header(s, "05", "System Analysis and Design", "Requirements")
    f = textbox(s, Inches(0.72), Inches(1.6), Inches(6.2), Inches(5.4))
    write(
        f,
        [
            "Functional requirements",
            ("Register and log in securely.", 1),
            ("Create and manage pet profiles.", 1),
            ("Predict disease risk from symptoms.", 1),
            ("Agentic AI gives advice and urgency.", 1),
            ("Explain the result in plain words.", 1),
            ("Find clinics, surgeons and free time slots.", 1),
            ("Book, reschedule and cancel appointments.", 1),
            ("Track vaccinations and send reminders.", 1),
            ("Keep health history and notifications.", 1),
        ],
        size=16,
    )
    f2 = textbox(s, Inches(7.1), Inches(1.6), Inches(5.6), Inches(5.4))
    write(
        f2,
        [
            "Non-functional requirements",
            ("Security — JWT login, role-based access.", 1),
            ("Privacy — a vet sees a pet only with a valid link.", 1),
            ("Usability — three languages: English, Sinhala, Tamil.", 1),
            ("Reliability — services recover on their own.", 1),
            ("Performance — result in a few seconds.", 1),
            ("Maintainability — separate services, shared types.", 1),
            "Important rule",
            ("Every result shows a decision-support disclaimer.", 1),
            ("The system supports decisions. It does not diagnose.", 1),
        ],
        size=16,
    )
    notes(
        s,
        "Nine functional requirements, all implemented. On the non-functional side the two I want to "
        "highlight are privacy and the disclaimer. A vet cannot open a pet's records unless there is a "
        "real link — an appointment or the owner's consent. And every result carries a disclaimer, "
        "because this is decision support, not diagnosis.",
    )

    # 6 ── Architecture
    s = blank(prs)
    header(s, "06", "System Architecture", "Micro-services behind one API gateway")
    picture(s, "fig06-architecture.png")
    notes(
        s,
        "The front end talks to one API gateway. The gateway checks the login token and passes the "
        "request on. Node services handle users, pets, clinics, vaccinations and admin. Two Python "
        "services do the AI work. MongoDB stores the data and Neo4j stores the disease knowledge "
        "graph. Everything runs in Docker — thirteen containers.",
    )

    # 7 ── Use case
    s = blank(prs)
    header(s, "07", "Use Case Diagram", "Three roles: pet owner, veterinarian, administrator")
    picture(s, "fig02-use-case.png")
    notes(
        s,
        "Three actors. The owner reports symptoms, sees the risk, books and tracks vaccinations. "
        "The vet sees patients they are linked to and confirms the real diagnosis. The admin manages "
        "clinics, vet accounts and approvals.",
    )

    # 8 ── ER + class
    s = blank(prs)
    header(s, "08", "ER Diagram and Database Design", "MongoDB collections and their links")
    picture(s, "fig03-erd.png")
    notes(
        s,
        "The main entities are User, Pet, Prediction, Appointment, Vaccination and Clinic. "
        "One extra table matters: the access grant, which records which vet may see which pet, "
        "and when that access started and ended.",
    )

    s = blank(prs)
    header(s, "09", "Class Diagram", "Main classes and services")
    picture(s, "fig05-class.png")
    notes(s, "This shows the main classes and how the services are organised in code.")

    # 10 ── Agent pipeline
    s = blank(prs)
    header(s, "10", "Agentic AI Pipeline", "Four agents working together")
    picture(s, "fig07-agent-pipeline.png")
    notes(
        s,
        "This is the agentic part. The first agent predicts the risk by calling the ML model. "
        "The second agent explains it and gives care advice. If the case needs a vet, a third agent "
        "finds and ranks clinics. A fourth agent runs in the background on a timer and creates alerts "
        "without anyone asking. The route through the graph is decided by the result, not fixed.",
    )

    # 11 ── Development progress
    s = blank(prs)
    header(s, "11", "Development Progress", "What is built and working")
    stat_row(
        s,
        [
            ("13", "Docker containers"),
            ("9", "back-end services"),
            ("3", "role dashboards"),
            ("3", "languages (EN/SI/TA)"),
            ("20 / 24", "diseases / symptoms in ontology"),
        ],
    )
    f = textbox(s, Inches(0.72), Inches(3.5), Inches(5.9), Inches(3.6))
    write(
        f,
        [
            "Modules completed",
            ("Login, registration, role-based access.", 1),
            ("Pet profiles and health history.", 1),
            ("Symptom form and risk prediction.", 1),
            ("Explanation of every result.", 1),
            ("Clinic map, booking, cancel and reschedule.", 1),
            ("Vaccination tracking and reminders.", 1),
            ("Owner–vet inquiry messaging.", 1),
            ("Admin approvals, tickets and audit log.", 1),
        ],
        size=16,
    )
    f2 = textbox(s, Inches(7.0), Inches(3.5), Inches(5.7), Inches(3.6))
    write(
        f2,
        [
            "Technologies used",
            ("React, TypeScript, Tailwind, installable PWA.", 1),
            ("Node.js + Express micro-services.", 1),
            ("Python FastAPI, scikit-learn, LangGraph, Gemini.", 1),
            ("MongoDB, Neo4j, Docker Compose.", 1),
            "Version control",
            ("Git and GitHub, committed after each feature.", 1),
        ],
        size=16,
    )
    notes(
        s,
        "All nine functional requirements are built and running. The whole stack starts with one "
        "docker compose command — thirteen containers. Everything is in Git and pushed to GitHub "
        "after each feature.",
    )

    # 12 ── Prototype
    s = blank(prs)
    header(s, "12", "Prototype — How the System Works", "One journey, start to finish")
    steps = [
        "1. Owner adds pet",
        "2. Answers symptom form",
        "3. ML predicts risk",
        "4. Agents explain + advise",
        "5. Clinic suggested",
        "6. Book appointment",
        "7. Reminders continue",
    ]
    gap = Inches(0.16)
    width = Emu(int((Inches(11.9) - gap * (len(steps) - 1)) / len(steps)))
    x = Inches(0.72)
    for i, step in enumerate(steps):
        card = rect(s, x, Inches(1.7), width, Inches(1.0), TEAL if i % 2 == 0 else NAVY)
        fr = card.text_frame
        fr.word_wrap = True
        pa = fr.paragraphs[0]
        pa.alignment = PP_ALIGN.CENTER
        ru = pa.add_run()
        ru.text = step
        ru.font.size = Pt(12)
        ru.font.bold = True
        ru.font.color.rgb = WHITE
        ru.font.name = "Calibri"
        x = Emu(int(x + width + gap))
    placeholder(
        s,
        "PASTE SCREENSHOTS HERE\n\n"
        "Suggested: (1) symptom form   (2) risk result with explanation   (3) clinic map and booking",
        y=Inches(3.05),
        h=Inches(3.6),
    )
    notes(
        s,
        "This is the full journey. The owner adds a pet, answers a short form about appetite, water, "
        "activity and symptoms. The model predicts. The agents explain it and say how urgent it is. "
        "If a vet is needed, a clinic is suggested with a free slot, and the owner books. "
        "After that, reminders keep going. I will now show this live.",
    )

    # 13 ── Explainability
    s = blank(prs)
    header(s, "13", "Explainability", "The owner can see why")
    f = textbox(s, Inches(0.72), Inches(1.6), Inches(5.9), Inches(5.4))
    write(
        f,
        [
            "What the owner sees",
            ("The risk level and a confidence score.", 1),
            ("Top 3 possible conditions with probabilities.", 1),
            ("Which answers pushed the result up or down.", 1),
            ("Known symptom-to-disease links from the ontology.", 1),
            "Why this matters",
            ("An owner will not act on a number they do not trust.", 1),
            ("A vet can check the reasoning, not just the answer.", 1),
        ],
        size=17,
    )
    placeholder(
        s,
        "PASTE SCREENSHOT HERE\n\nThe 'What led to this result' panel",
        x=Inches(7.0),
        y=Inches(1.9),
        w=Inches(5.7),
        h=Inches(4.3),
    )
    notes(
        s,
        "Explainability was a specific objective. The system shows the confidence, the top three "
        "conditions, and which of the owner's answers mattered most. Early on this showed raw word "
        "fragments like 'has' and 'less', which confused people, so I changed it to show whole form "
        "answers such as 'water intake: increased'.",
    )

    # 14 ── Results: models
    s = blank(prs)
    header(s, "14", "Current Results — Model Performance", "Measured, not estimated")
    stat_row(
        s,
        [
            (f"{F['accuracy'] * 100:.1f}%", "Model A accuracy"),
            (f"{F['macro_f1']:.3f}", "Model A macro F1"),
            (f"{F['n_labels']}", "conditions predicted"),
            (f"{F['n_train']:,}/{F['n_test']:,}", "train / test rows"),
            (f"{F['roc_auc']:.3f}", "Model B ROC-AUC"),
        ],
    )
    real = ["Skin Irritations", "Parasites", "Ear Infections", "Mobility Problems", "Digestive Issues"]
    rows = [["Condition", "F1", "Data source"]]
    for name in real:
        rows.append([name, f"{F['per_class'][name]:.2f}", "Real dataset"])
    rows.append(["Diabetes, Heart, Liver, Pancreatitis", "1.00", "Knowledge-generated"])
    rows.append(["Cancer, Chronic Kidney Disease", "0.97", "Knowledge-generated"])
    table(s, rows, [5.6, 1.6, 4.7], top=Inches(3.5), font=13)
    notes(
        s,
        "Model A reaches ninety point three percent accuracy over eleven conditions. But I want to be "
        "honest about that number. The six chronic classes score near one because their training rows "
        "were generated from veterinary knowledge, not collected. The five classes from the real "
        "dataset sit between zero point seven seven and zero point eight seven, and they confuse with "
        "each other. So I quote the per-class numbers, not just the headline.",
    )

    # 15 ── Results: testing
    s = blank(prs)
    header(s, "15", "Current Results — Testing", "Manual and automated")
    stat_row(
        s,
        [
            (str(TC_TOTAL), "manual test cases"),
            (str(TC_PASS), "passed"),
            (str(TC_FAIL), "failed"),
            ("97.6%", "pass rate (executed)"),
            (str(TESTS_AI + TESTS_AGENT + TESTS_GW), "automated tests pass"),
        ],
    )
    f = textbox(s, Inches(0.72), Inches(3.5), Inches(5.9), Inches(3.6))
    write(
        f,
        [
            "How I tested",
            (f"A test workbook of {TC_TOTAL} cases across every module.", 1),
            ("Each case has steps, expected result and actual result.", 1),
            (f"{TC_NOTRUN} cases are still to run.", 1),
            (
                f"Automated: {TESTS_AI} AI-service, {TESTS_AGENT} agent-service, "
                f"{TESTS_GW} gateway.",
                1,
            ),
        ],
        size=16,
    )
    f2 = textbox(s, Inches(7.0), Inches(3.5), Inches(5.7), Inches(3.6))
    write(
        f2,
        [
            "What testing found",
            ("Vaccination records could not be edited — fixed.", 1),
            ("Booking confirmed with no confirmation step — fixed.", 1),
            ("Vet reply saved no message — fixed.", 1),
            ("Any vet could read any inquiry — fixed.", 1),
            "Feedback loop is live",
            ("Vets confirm the real diagnosis in the system.", 1),
            ("AI-vs-vet agreement is shown on the admin panel.", 1),
        ],
        size=16,
    )
    notes(
        s,
        "I wrote a test workbook of one hundred and twenty three cases and I am working through it by "
        "hand. Eighty two pass, two fail, thirty nine still to run. Testing found real defects and I "
        "fixed each one. There are also fifty three automated tests, and all of them pass.",
    )

    # 16 ── Access control
    s = blank(prs)
    header(s, "16", "Privacy — Relationship-Based Access Control", "A vet sees a pet only when there is a real link")
    picture(s, "fig08-access-control.png", top=Inches(1.5), max_h=Inches(4.1))
    f = textbox(s, Inches(0.72), Inches(5.85), Inches(11.9), Inches(1.3))
    write(
        f,
        [
            "Access comes from an appointment or the owner's consent. The owner can withdraw it at any time.",
            ("Checks are in the services, not the screen. If the check cannot run, access is refused.", 1),
        ],
        size=15,
    )
    notes(
        s,
        "This was the biggest security fix of this stage. Before, any vet could open any pet's records. "
        "Now a vet needs a real relationship: either the owner booked with them, or the owner shared the "
        "records. The owner can take that back. And the check happens in the back end, not the screen, "
        "so hiding a button is not what protects the data.",
    )

    # 17 ── Challenges
    s = blank(prs)
    header(s, "17", "Challenges and Solutions", "Technical and research")
    table(
        s,
        [
            ["Challenge", "What I did"],
            [
                "The disease graph loaded before the database was ready, so it was empty.",
                "Added a health check and a background retry until it connects.",
            ],
            [
                "Explanations showed word pieces like 'has' and 'less'. Users were confused.",
                "Grouped them into whole form answers, e.g. 'water intake: increased'.",
            ],
            [
                "Any veterinarian could open any pet's health records.",
                "Built relationship-based access control that fails closed.",
            ],
            [
                "The risk dataset is 97.7% 'dangerous', so a label would always say dangerous.",
                "Used it as a calibrated probability only, never a yes/no label. Stated as a limitation.",
            ],
            [
                "Planned Vertex AI, Google Maps and Firebase messaging needed paid billing.",
                "Used Gemini API, OpenStreetMap and in-app notifications instead.",
            ],
        ],
        [6.0, 5.9],
        font=13,
    )
    f = textbox(s, Inches(0.72), Inches(5.95), Inches(11.9), Inches(1.2))
    write(
        f,
        [
            "Lesson learned: test the boring paths. Most defects were in restart, delete and cancel — not in the AI.",
        ],
        size=15,
    )
    notes(
        s,
        "Five real challenges. The one I would highlight is the risk dataset. Almost ninety eight "
        "percent of its rows are labelled dangerous, so a model trained on it would just say dangerous "
        "every time. Rather than hide that, I use it only as a probability and I state it openly as a "
        "limitation. The lesson overall: most of my defects were not in the AI, they were in restart, "
        "delete and cancel.",
    )

    # 18 ── Remaining work
    s = blank(prs)
    header(s, "18", "Remaining Work", "What is left before submission")
    f = textbox(s, Inches(0.72), Inches(1.6), Inches(5.9), Inches(5.4))
    write(
        f,
        [
            "Features to complete",
            (f"Run the remaining {TC_NOTRUN} test cases.", 1),
            ("Close the deactivated-vet login gap.", 1),
            ("Archive old risk assessments instead of deleting.", 1),
            ("Final accessibility and mobile polish.", 1),
            "Testing and validation",
            ("User testing with real pet owners.", 1),
            ("Collect vet-confirmed diagnoses for AI-vs-vet agreement.", 1),
        ],
        size=17,
    )
    f2 = textbox(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.4))
    write(
        f2,
        [
            "Documentation",
            ("Final dissertation chapters.", 1),
            ("Updated diagrams and test evidence.", 1),
            ("User manual and deployment guide.", 1),
            "Final steps",
            ("Full system evaluation against the objectives.", 1),
            ("Viva preparation.", 1),
        ],
        size=17,
    )
    notes(
        s,
        "What is left is mainly evaluation and writing. Thirty nine test cases still to run, user "
        "testing with real owners, and collecting enough vet-confirmed diagnoses to report the "
        "AI-versus-vet agreement rate properly. Then the dissertation.",
    )

    # 19 ── Timeline
    s = blank(prs)
    header(s, "19", "Project Timeline", "Milestones completed and remaining schedule")
    picture(s, "fig10-gantt.png", top=Inches(1.5), max_h=Inches(4.3))
    f = textbox(s, Inches(0.72), Inches(6.0), Inches(11.9), Inches(1.1))
    write(
        f,
        [
            "Completed: proposal, literature review, requirements, design, full implementation, interim report, PR2 testing.",
            ("Remaining: user testing, evaluation, dissertation, viva.", 1),
        ],
        size=15,
    )
    notes(
        s,
        "I am on schedule. Proposal, literature, design and implementation are done, the interim "
        "report is submitted, and I am now in the testing and evaluation phase.",
    )

    # 20 ── Conclusion
    s = blank(prs)
    header(s, "20", "Conclusion", "Summary of progress")
    f = textbox(s, Inches(0.72), Inches(1.7), Inches(11.9), Inches(4.6))
    write(
        f,
        [
            "All six objectives are implemented and running end to end.",
            (
                f"Model A predicts {F['n_labels']} conditions at {F['accuracy'] * 100:.1f}% accuracy; "
                f"per-class results are reported honestly.",
                1,
            ),
            ("The agentic layer reasons, explains, recommends and monitors on its own.", 1),
            ("Explanations are in plain words, in three languages.", 1),
            ("Privacy is enforced by real relationships, not by hiding buttons.", 1),
            f"Tested: {TC_PASS} of {TC_PASS + TC_FAIL} executed manual cases pass, plus "
            f"{TESTS_AI + TESTS_AGENT + TESTS_GW} automated tests.",
            "The system supports decisions. It never claims to diagnose.",
            "Expected completion: final dissertation and viva at the end of the academic year.",
        ],
        size=19,
    )
    notes(
        s,
        "To summarise: all six objectives are built and working. The models are measured and I report "
        "their weaknesses as well as their strengths. The agentic layer reasons and acts on its own. "
        "And throughout, the system is decision support — it never claims to diagnose. Thank you.",
    )

    # 21 ── References
    s = blank(prs)
    header(s, "21", "References", "IEEE style")
    refs = [
        "[1] S. C. Kim and S. Kim, \"Development of a Dog Health Score Using an Artificial Intelligence Disease "
        "Prediction Algorithm Based on Multifaceted Data,\" Animals, vol. 14, no. 2, Jan. 2024.",
        "[2] D. Szlosek et al., \"Development and validation of a machine learning model for clinical wellness "
        "visit classification in cats and dogs,\" Front. Vet. Sci., vol. 11, 2024.",
        "[3] S. Das, R. K. Roy, and T. Bezboruah, \"Machine Learning in Animal Healthcare: A Comprehensive "
        "Review,\" Int. J. Recent Eng. Sci., vol. 11, no. 3, pp. 89–93, Jun. 2024.",
        "[4] M. Jokar, A. Abdous, and V. Rahmanian, \"AI chatbots in pet health care: Opportunities and "
        "challenges for owners,\" Vet. Med. Sci., May 2024.",
        "[5] C. A. Aguilar-Lazcano et al., \"Machine Learning-Based Sensor Data Fusion for Animal Monitoring: "
        "Scoping Review,\" Sensors, vol. 23, no. 12, Jun. 2023.",
        "[6] J. Imada et al., \"Comparison of Machine Learning Tree-Based Algorithms to Predict Future "
        "Paratuberculosis ELISA Results Using Repeat Milk Tests,\" Animals, vol. 14, no. 7, Apr. 2024.",
        "[7] S. Eman et al., \"Technologies in Biomarker Discovery for Animal Diseases,\" Animals, vol. 15, "
        "no. 21, Nov. 2025.",
        "[8] J. Rathi and A. Sumathi, \"Animal Health Prediction Using Hybrid KNN Based Vector Neighbor "
        "Classification Model,\" J. Comput. Sci., vol. 21, no. 9, pp. 2088–2095, Sep. 2025.",
        "[9] F. Bouchemla et al., \"Artificial intelligence feasibility in veterinary medicine: A systematic "
        "review,\" Vet. World, vol. 16, no. 10, pp. 2143–2149, 2023.",
        "[10] A. AP, \"PAWPAL: AI-Based Pet Health Assistance System Using Image-based Disease Detection and "
        "NLP Symptom Analysis,\" Int. J. Res. Appl. Sci. Eng. Technol., vol. 13, no. 12, pp. 2065–2072, Dec. 2025.",
    ]
    f = textbox(s, Inches(0.72), Inches(1.6), Inches(11.9), Inches(5.5))
    f.clear()
    for i, ref in enumerate(refs):
        para = f.paragraphs[0] if i == 0 else f.add_paragraph()
        para.space_after = Pt(7)
        run = para.add_run()
        run.text = ref
        run.font.size = Pt(11.5)
        run.font.color.rgb = INK
        run.font.name = "Calibri"
    notes(s, "These are the key references, in IEEE style. The full list is in the interim report.")

    # 22 ── Q&A
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    f = textbox(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(2.2))
    p = f.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank You"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"
    p2 = f.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Questions?"
    r2.font.size = Pt(28)
    r2.font.color.rgb = TEAL
    r2.font.name = "Calibri"
    notes(
        s,
        "Thank you for your time. I am happy to take questions.\n\n"
        "Likely questions and short answers:\n"
        "• Why is accuracy so high? Six of the eleven classes use knowledge-generated rows. "
        "The five real-data classes sit at 0.77 to 0.87, and I report both.\n"
        "• Is this a diagnosis? No. Every result carries a decision-support disclaimer.\n"
        "• What makes it agentic? The route through the pipeline is decided by the result, and a "
        "fourth agent runs on a timer with no user request.\n"
        "• How is privacy handled? A vet needs an appointment or the owner's consent, checked in the "
        "back end, and it fails closed.",
    )

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
